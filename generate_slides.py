from __future__ import annotations

import re
from html import unescape
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
SLIDES = ROOT / "slides"

COWORK_TASKS = {
    "Light": {
        "All teams": [
            ("Rebalance your week and protect focus time", "Rebalance the week, resolve conflicts, and protect focus blocks."),
            ("Wrap up projects and organize all related work", "Close the project with a clean, shareable OneDrive archive."),
            ("Catch up on messages and send replies automatically", "Handle routine replies and draft higher-stakes responses for review."),
        ],
        "Sales": [
            ("Pricing screenshot to customer presentation", "Turn approved pricing into a customer-ready deck and draft email."),
            ("Weekly external customer email review", "Surface action items, draft replies, and send a Teams reminder."),
        ],
        "Marketing": [
            ("Research and insights alignment recap", "Post decisions, open issues, and named next steps to Teams."),
        ],
    },
    "Medium": {
        "All teams": [
            ("Turn inbox noise into a curated intelligence brief", "Create a weekly trend brief from the publications you follow."),
        ],
        "Sales": [],
        "Marketing": [
            ("Partner and channel activation kit", "Produce partner-ready and channel-ready kits with a routing plan."),
        ],
    },
    "Heavy": {
        "All teams": [
            ("Prepare a complete out-of-office handoff", "Build a handoff from projects, owners, deadlines, and files."),
            ("Analyze and optimize your OneDrive at scale", "Create a catalog with cleanup decisions queued for approval."),
            ("Audit and surface against a standard", "Audit files against a policy and return a sortable fix list."),
            ("Onboard a new hire with a complete 30-60-90 plan", "Create the plan, getting started dashboard, and check-ins."),
        ],
        "Sales": [
            ("Product launch customer pitch", "Create the comparison, value prop, and customer-ready pitch deck."),
            ("New customer onboarding automation", "Set up onboarding artifacts, team message, and welcome email."),
            ("Customer adoption materials", "Build persona-tailored learning paths and rollout order."),
            ("ROI and value selling artifact", "Produce an executive deck and microsite with use cases and ROI model."),
        ],
        "Marketing": [
            ("Analyst briefing prep and rehearsal routing", "Build a briefing dashboard and lock rehearsal time."),
            ("Messaging drift audit and remediation routing", "Flag inconsistencies with severity, fix, and owner."),
            ("Scheduled customer signal activation", "Create a Friday brief with themes and campaign adjustments."),
        ],
    },
}

SCHEDULED_PROMPT_TIP = (
    'Any of these Cowork tasks can be scheduled to run automatically. '
    'In Microsoft 365 Copilot, hover a saved prompt and pick "Schedule this prompt", '
    'then choose a daily or weekly cadence (requires a Microsoft 365 Copilot licence).'
)

PAGES = [
    ("business.html", "What to use when - Copilot, Cowork, Studio and Scout (Business - full).pptx", "business"),
    ("business-lite.html", "What to use when - Copilot and Cowork (Business - lite).pptx", "business"),
    ("technical.html", "What to use when - Copilot, Cowork, Studio and Scout (Technical - full).pptx", "technical"),
    ("technical-lite.html", "What to use when - Copilot and Cowork (Technical - lite).pptx", "technical"),
]

THEMES = {
    "business": {
        "bg": RGBColor(255, 250, 245),
        "panel": RGBColor(255, 255, 255),
        "ink": RGBColor(27, 27, 44),
        "muted": RGBColor(90, 96, 120),
        "accent": RGBColor(26, 111, 214),
        "accent2": RGBColor(62, 142, 34),
        "accent3": RGBColor(199, 62, 120),
        "soft": RGBColor(246, 247, 251),
    },
    "technical": {
        "bg": RGBColor(14, 20, 48),
        "panel": RGBColor(27, 34, 65),
        "ink": RGBColor(245, 247, 255),
        "muted": RGBColor(174, 184, 224),
        "accent": RGBColor(0, 183, 195),
        "accent2": RGBColor(115, 184, 46),
        "accent3": RGBColor(227, 80, 143),
        "soft": RGBColor(22, 30, 69),
    },
}

TIER_COLORS = {
    "Light": RGBColor(62, 142, 34),
    "Medium": RGBColor(14, 142, 146),
    "Heavy": RGBColor(199, 62, 120),
}

SECTION_STOP_WORDS = {"System", "Light", "Dark", "References"}


def clean_text(value: str) -> str:
    value = re.sub(r"<script[\s\S]*?</script>", " ", value, flags=re.I)
    value = re.sub(r"<style[\s\S]*?</style>", " ", value, flags=re.I)
    value = re.sub(r"<[^>]+>", " ", value)
    value = unescape(value)
    value = (
        value.replace("\u2013", ",")
        .replace("\u2014", ",")
        .replace("\u2018", "'")
        .replace("\u2019", "'")
        .replace("\u201c", '"')
        .replace("\u201d", '"')
    )
    value = re.sub(r"\s+", " ", value).strip()
    return value


def find_first(pattern: str, text: str) -> str:
    match = re.search(pattern, text, flags=re.I | re.S)
    return clean_text(match.group(1)) if match else ""


def split_sections(html: str) -> list[str]:
    return re.findall(r"<section[\s\S]*?</section>", html, flags=re.I)


def extract_bullets(section_html: str, title: str, limit: int = 7) -> list[str]:
    candidates = []
    selectors = [
        r'<div class="nm">([\s\S]*?)</div>',
        r'<div class="task">([\s\S]*?)</div>',
        r'<div class="t">([\s\S]*?)</div>',
        r'<div class="title">([\s\S]*?)</div>',
        r'<h3>([\s\S]*?)</h3>',
        r'<div class="lab">([\s\S]*?)</div>\s*<h3>([\s\S]*?)</h3>',
        r'<td class="task">([\s\S]*?)</td>',
    ]
    for pattern in selectors:
        for match in re.finditer(pattern, section_html, flags=re.I | re.S):
            if len(match.groups()) == 2:
                text = clean_text(match.group(1) + ": " + match.group(2))
            else:
                text = clean_text(match.group(1))
            if text and text not in candidates and text not in SECTION_STOP_WORDS and text != title:
                candidates.append(text)
    if len(candidates) < 3:
        for match in re.finditer(r"<p[^>]*>([\s\S]*?)</p>", section_html, flags=re.I | re.S):
            text = clean_text(match.group(1))
            if 20 <= len(text) <= 170 and text not in candidates:
                candidates.append(text)
    return candidates[:limit]


def parse_page(path: Path) -> dict:
    html = path.read_text(encoding="utf-8")
    title = find_first(r"<h1[^>]*>([\s\S]*?)</h1>", html) or clean_text(find_first(r"<title>([\s\S]*?)</title>", html))
    lede = find_first(r'<p class="lede">([\s\S]*?)</p>', html) or find_first(r'<div class="top-sub">([\s\S]*?)</div>', html)
    sections = []
    for section in split_sections(html):
        heading = find_first(r"<h2[^>]*>([\s\S]*?)</h2>", section)
        if not heading:
            continue
        if heading == "Cowork tasks in action":
            sections.append({"title": heading, "special": "cowork"})
            continue
        subtitle = find_first(r'<p class="sec-sub">([\s\S]*?)</p>', section) or find_first(r'<span class="hint">([\s\S]*?)</span>', section)
        bullets = extract_bullets(section, heading)
        sections.append({"title": heading, "subtitle": subtitle, "bullets": bullets, "special": None})
    refs_raw = re.findall(r'<a href="([^"]+)"[^>]*>([\s\S]*?)</a>', html, flags=re.I | re.S)
    seen = set()
    refs = []
    for url, label in refs_raw:
        label_clean = clean_text(label)
        if not label_clean or label_clean in SECTION_STOP_WORDS:
            continue
        if not url.startswith("http"):
            continue
        key = (url, label_clean)
        if key in seen:
            continue
        seen.add(key)
        refs.append({"url": url, "label": label_clean})
    return {"title": title, "lede": lede, "sections": sections, "refs": refs[:8]}


def extract_label_value(detail_html: str, label: str) -> str:
    match = re.search(rf"<p><b>{re.escape(label)}:</b>\s*([\s\S]*?)</p>", detail_html, flags=re.I)
    return clean_text(match.group(1)) if match else ""


def extract_cowork_task_details(path: Path) -> list[dict]:
    html = path.read_text(encoding="utf-8")
    tasks = []
    detail_blocks = [
        match.group(1)
        for match in re.finditer(r'<details class="task-detail">([\s\S]*?)</details>', html, flags=re.I)
    ]
    if not detail_blocks:
        starts = list(re.finditer(r'<div class="task-modal-data" id="[^"]+">', html, flags=re.I))
        detail_blocks = [
            html[match.start() : starts[index + 1].start() if index + 1 < len(starts) else len(html)]
            for index, match in enumerate(starts)
        ]
    for detail in detail_blocks:
        tier = find_first(r'<span class="tier-badge\s+([^"]+)">', detail).title()
        title = find_first(r'<span class="task-title">([\s\S]*?)</span>', detail)
        prompt_match = re.search(r'<pre class="prompt-text">([\s\S]*?)</pre>', detail, flags=re.I)
        prompt = unescape(prompt_match.group(1)).strip() if prompt_match else ""
        prompt = clean_text(prompt)
        flow_match = re.search(r'<ol class="exec-flow">([\s\S]*?)</ol>', detail, flags=re.I)
        flow = [clean_text(item) for item in re.findall(r"<li>([\s\S]*?)</li>", flow_match.group(1), flags=re.I)] if flow_match else []
        notes_match = re.search(r'<ul class="task-notes">([\s\S]*?)</ul>', detail, flags=re.I)
        notes = [clean_text(item) for item in re.findall(r"<li>([\s\S]*?)</li>", notes_match.group(1), flags=re.I)] if notes_match else []
        tasks.append(
            {
                "tier": tier,
                "title": title,
                "why": extract_label_value(detail, f"Why {tier.lower()}"),
                "goal": extract_label_value(detail, "Goal"),
                "outcome": extract_label_value(detail, "Outcome"),
                "prompt": prompt,
                "flow": flow,
                "notes": notes,
            }
        )
    return tasks


def set_background(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, size: int, color: RGBColor, bold: bool = False, align=None, anchor=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items: Iterable[str], x: float, y: float, w: float, h: float, theme: dict, size: int = 14) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = theme["ink"]
        p.space_after = Pt(7)
        p.line_spacing = 1.15


def style_title_placeholder(slide, text: str, theme: dict, size: int = 28):
    title_shape = slide.shapes.title
    title_shape.left = Inches(0.7)
    title_shape.top = Inches(0.55)
    title_shape.width = Inches(12.0)
    title_shape.height = Inches(0.9)
    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = "Segoe UI"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = theme["ink"]


def add_header(slide, title: str, theme: dict, kicker: str | None = None) -> None:
    style_title_placeholder(slide, title, theme, 28)
    if kicker:
        add_textbox(slide, kicker.upper(), 0.7, 0.28, 10.0, 0.28, 11, theme["muted"], True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.5), Inches(1.6), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = theme["accent"]
    line.line.fill.background()


def add_panel(slide, x: float, y: float, w: float, h: float, theme: dict, accent: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["panel"]
    shape.line.color.rgb = theme["soft"]
    shape.line.width = Pt(0.75)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.09))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return shape


def add_title_slide(prs: Presentation, page: dict, theme: dict, kind: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide, theme["bg"])
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(-0.6), Inches(2.8), Inches(2.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme["accent"]
    accent.line.fill.background()
    title_shape = slide.shapes.title
    title_shape.left = Inches(0.9)
    title_shape.top = Inches(2.0)
    title_shape.width = Inches(10.5)
    title_shape.height = Inches(1.8)
    title_shape.text_frame.text = page["title"]
    title_shape.text_frame.word_wrap = True
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.name = "Segoe UI"
        run.font.size = Pt(38)
        run.font.bold = True
        run.font.color.rgb = theme["ink"]
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.left = Inches(0.92)
            ph.top = Inches(4.05)
            ph.width = Inches(11.0)
            ph.height = Inches(1.6)
            ph.text_frame.text = page["lede"]
            ph.text_frame.word_wrap = True
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.name = "Segoe UI"
                run.font.size = Pt(18)
                run.font.color.rgb = theme["muted"]
    label = "Warm customer guide" if kind == "business" else "Technical decision guide"
    add_textbox(slide, label, 0.92, 6.45, 6.0, 0.4, 14, theme["accent"], True)


def add_section_slide(prs: Presentation, section: dict, theme: dict, idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, section["title"], theme, f"Section {idx:02d}")
    content_top = 1.75
    if section.get("subtitle"):
        add_textbox(slide, section["subtitle"], 0.7, 1.7, 12.0, 0.55, 15, theme["muted"])
        content_top = 2.4
    bullets = section.get("bullets") or ["Use this page section as the decision checkpoint for the workflow."]
    left = bullets[:4]
    right = bullets[4:8]
    panel_h = 7.5 - content_top - 0.5
    add_panel(slide, 0.7, content_top, 5.95, panel_h, theme, theme["accent"])
    add_bullets(slide, left, 0.95, content_top + 0.25, 5.55, panel_h - 0.4, theme, 14)
    add_panel(slide, 6.85, content_top, 5.95, panel_h, theme, theme["accent2"])
    if right:
        add_bullets(slide, right, 7.1, content_top + 0.25, 5.55, panel_h - 0.4, theme, 14)
    else:
        add_textbox(slide, "Key takeaway", 7.1, content_top + 0.3, 5.5, 0.4, 17, theme["ink"], True)
        add_textbox(slide, "Pick the simplest Copilot surface that can finish the job with the right level of control.", 7.1, content_top + 0.85, 5.5, 1.6, 15, theme["muted"])


def add_cowork_summary(prs: Presentation, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, "Cowork tasks: light, medium, heavy", theme, "Copilot Cowork")
    notes = [
        ("Light", "Quick everyday handoffs, minutes to finish, a source or two."),
        ("Medium", "Recurring multi-step work that reasons across several sources."),
        ("Heavy", "Large multi-part jobs that produce a full set of artifacts."),
    ]
    card_w = 3.85
    gap = 0.25
    total = 3 * card_w + 2 * gap
    start_x = (13.333 - total) / 2
    for i, (tier, body) in enumerate(notes):
        x = start_x + i * (card_w + gap)
        add_panel(slide, x, 2.05, card_w, 3.7, theme, TIER_COLORS[tier])
        add_textbox(slide, tier, x, 2.4, card_w, 0.6, 28, TIER_COLORS[tier], True, PP_ALIGN.CENTER)
        add_textbox(slide, body, x + 0.3, 3.2, card_w - 0.6, 2.0, 15, theme["muted"], False, PP_ALIGN.CENTER)
    add_textbox(slide, "Classification follows the source task packs: Lightweight, Everyday workflows, and Hard Problems.", 0.9, 6.45, 11.5, 0.4, 12, theme["muted"], False, PP_ALIGN.CENTER)


def add_cowork_tier_slide(prs: Presentation, tier: str, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, f"Cowork tasks: {tier.lower()}", theme, "Grouped by audience")
    audiences = [(k, v) for k, v in COWORK_TASKS[tier].items() if v]
    gap = 0.25
    total_w = 12.6
    col_w = (total_w - gap * (len(audiences) - 1)) / len(audiences)
    top = 1.85
    panel_h = 5.15
    for i, (audience, items) in enumerate(audiences):
        x = 0.7 + i * (col_w + gap)
        add_panel(slide, x, top, col_w, panel_h, theme, TIER_COLORS[tier])
        add_textbox(slide, audience, x + 0.2, top + 0.2, col_w - 0.4, 0.4, 16, TIER_COLORS[tier], True)
        y = top + 0.8
        item_count = len(items)
        if item_count <= 2:
            title_size, outcome_size, row_h, title_h = 14, 12, 1.5, 0.7
        elif item_count == 3:
            title_size, outcome_size, row_h, title_h = 13, 11, 1.25, 0.65
        else:
            title_size, outcome_size, row_h, title_h = 12, 10, 1.05, 0.6
        for title, outcome in items:
            add_textbox(slide, title, x + 0.2, y, col_w - 0.4, title_h, title_size, theme["ink"], True)
            add_textbox(slide, outcome, x + 0.2, y + title_h, col_w - 0.4, row_h - title_h - 0.05, outcome_size, theme["muted"])
            y += row_h


def compact_text(text: str, limit: int = 180) -> str:
    text = clean_text(text)
    if len(text) <= limit:
        return text
    return text[: limit - 1].rsplit(" ", 1)[0] + "..."


def select_featured_cowork_tasks(tasks: list[dict]) -> list[dict]:
    featured = []
    for tier in ["Light", "Medium", "Heavy"]:
        tier_tasks = [task for task in tasks if task["tier"] == tier]
        featured.extend(tier_tasks[:2])
    return featured


def add_cowork_task_overview(prs: Presentation, task: dict, theme: dict) -> None:
    tier = task["tier"]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, compact_text(task["title"], 80), theme, f"Cowork {tier} task")
    add_panel(slide, 0.7, 1.85, 12.0, 4.95, theme, TIER_COLORS[tier])
    add_textbox(slide, tier, 1.0, 2.2, 1.8, 0.55, 24, TIER_COLORS[tier], True)
    add_textbox(slide, compact_text(task["why"], 240), 2.95, 2.2, 9.4, 1.1, 15, theme["muted"])
    add_textbox(slide, "Goal", 1.0, 3.7, 1.7, 0.4, 16, theme["ink"], True)
    add_textbox(slide, compact_text(task["goal"], 280), 2.95, 3.7, 9.4, 1.2, 14, theme["ink"])
    add_textbox(slide, "Outcome", 1.0, 5.15, 1.85, 0.4, 16, theme["ink"], True)
    add_textbox(slide, compact_text(task["outcome"], 280), 2.95, 5.15, 9.4, 1.4, 14, theme["ink"])


def add_cowork_task_detail(prs: Presentation, task: dict, theme: dict) -> None:
    tier = task["tier"]
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, "Goal, outcome, prompt, flow", theme, compact_text(task["title"], 90))
    add_panel(slide, 0.7, 1.85, 5.95, 4.95, theme, TIER_COLORS[tier])
    add_textbox(slide, "Prompt to give Cowork", 0.95, 2.1, 5.5, 0.4, 15, theme["ink"], True)
    add_textbox(slide, compact_text(task["prompt"], 620), 0.95, 2.55, 5.5, 4.15, 11, theme["muted"])
    add_panel(slide, 6.85, 1.85, 5.95, 4.95, theme, TIER_COLORS[tier])
    add_textbox(slide, "Execution flow", 7.1, 2.1, 5.5, 0.4, 15, theme["ink"], True)
    flow = task["flow"][:6]
    has_notes = bool(task.get("notes"))
    flow_h = 3.0 if has_notes else 4.1
    add_bullets(slide, [compact_text(item, 120) for item in flow], 7.1, 2.55, 5.5, flow_h, theme, 12)
    if has_notes:
        note_text = " ".join(task["notes"])
        add_textbox(slide, "Note", 7.1, 5.7, 0.8, 0.3, 12, theme["ink"], True)
        add_textbox(slide, compact_text(note_text, 250), 7.1, 6.0, 5.5, 0.75, 11, theme["muted"])


def add_cowork_schedule_tip(prs: Presentation, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, "Add-on tip: schedule Cowork prompts", theme, "Copilot Cowork")
    add_panel(slide, 1.3, 2.4, 10.7, 3.4, theme, theme["accent2"])
    add_textbox(slide, "Tip", 1.7, 2.85, 1.6, 0.6, 28, theme["accent2"], True)
    add_textbox(slide, SCHEDULED_PROMPT_TIP, 1.7, 3.65, 9.9, 1.9, 19, theme["ink"])


def add_refs_slide(prs: Presentation, refs: list[dict], theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, theme["bg"])
    add_header(slide, "References", theme, "Microsoft sources")
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.85), Inches(11.5), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    for i, ref in enumerate(refs[:8]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.2
        bullet_run = p.add_run()
        bullet_run.text = "\u2022  "
        bullet_run.font.name = "Segoe UI"
        bullet_run.font.size = Pt(14)
        bullet_run.font.color.rgb = theme["ink"]
        link_run = p.add_run()
        link_run.text = ref["label"]
        link_run.font.name = "Segoe UI"
        link_run.font.size = Pt(14)
        link_run.font.color.rgb = theme["accent"]
        link_run.font.underline = True
        link_run.hyperlink.address = ref["url"]


def build_deck(page_file: str, out_file: str, kind: str) -> None:
    page_path = ROOT / page_file
    page = parse_page(page_path)
    cowork_task_details = extract_cowork_task_details(page_path)
    theme = THEMES[kind]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs, page, theme, kind)
    section_number = 1
    for section in page["sections"]:
        if section.get("special") == "cowork":
            add_cowork_summary(prs, theme)
            for tier in ["Light", "Medium", "Heavy"]:
                add_cowork_tier_slide(prs, tier, theme)
            for task in select_featured_cowork_tasks(cowork_task_details):
                add_cowork_task_overview(prs, task, theme)
                add_cowork_task_detail(prs, task, theme)
            add_cowork_schedule_tip(prs, theme)
        else:
            add_section_slide(prs, section, theme, section_number)
        section_number += 1
    add_refs_slide(prs, page["refs"], theme)
    SLIDES.mkdir(exist_ok=True)
    prs.save(SLIDES / out_file)
    print(f"Generated {SLIDES / out_file} with {len(prs.slides)} slides")


def main() -> None:
    for page_file, out_file, kind in PAGES:
        build_deck(page_file, out_file, kind)


if __name__ == "__main__":
    main()
